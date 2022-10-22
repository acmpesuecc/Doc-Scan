
from pdf2text import pdf2text
import cv2
import numpy as np
from comparejpg import compare_images
from tkinter import *

root = Tk()
root.geometry("600x800")
bg = PhotoImage(file="DOC-SCAN Tkinter UI.png")
label1 = Label(root, image=bg)
label1.place(x=0, y=0)
root.configure(bg="white")
import collections  # noqa
import collections.abc  # noqa
from pptx import Presentation  # noqa


def compare_images(img1, img2):
    a = cv2.imread(img1)
    b = cv2.imread(img2)
    difference = cv2.subtract(a, b)
    result = not np.any(difference)
    if result is True:
        print("Pictures are the same")
        Label(root, text=f"Pictures are the same", font=("Calibri", 15), bg="white").place(x=200, y=600)
    else:
        cv2.imwrite("difference.png", difference)
        Label(root, text=f"Pictures are different, the difference is stored as ed.png", font=("Calibri", 15), bg="white").pack()


def compare_ppt(p1, p2):

    ppt1 = Presentation(p1)
    ppt2 = Presentation(p2)

    text_runs1 = []
    text_runs2 = []

    for slide in ppt1.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs1.append(run.text)
    for slide in ppt2.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs2.append(run.text)
    if collections.Counter(text_runs1) == collections.Counter(text_runs2):
        Label(root, text=f"The presentations are same", font=("Calibri", 15), bg="white").place(x=200, y=600)

    else:
        Label(root, text=f"The presentations are different", font=("Calibri", 15), bg="white").place(x=200, y=600)
        Label(root, text=f"The differences are : ", font=("Calibri", 15), bg="white").place(x=200, y=630)
        Label(root, text=collections.Counter(text_runs1) - collections.Counter(text_runs2), font=("Calibri", 15), bg="white").place(x=200, y=640)


def comparetext(a, b):
    file_1 = open(a, 'r')
    file_2 = open(b, 'r')

    file_1_line = file_1.readline()
    file_2_line = file_2.readline()

    line_no = 1

    print()

    with open('text1.txt') as file1:
        with open('text2.txt') as file2:
            same = set(file1).intersection(file2)

    # print("Common Lines in Both Files")
    Label(root, text=f"Common Lines in Both Files", font=("Calibri", 15), bg="white").place(x=200, y=600)

    for line in same:
        # print(line, end='')
        Label(root, text=f"{line} ", font=("Calibri", 15), bg="white").pack()

    print('\n')
    Label(root, text=f"Difference Lines in Both Files", font=("Calibri", 15), bg="white").place(x=200, y=600)
    # print("Difference Lines in Both Files")
    print('\n')
    while file_1_line != '' or file_2_line != '':

        file_1_line = file_1_line.rstrip()
        file_2_line = file_2_line.rstrip()

        if file_1_line != file_2_line:

            if file_1_line == '':
                # print("file1-", "Line-%d" % line_no, file_1_line)
                Label(root, text=f"File1- {line_no} Line- {file_1_line}", font=("Calibri", 15), bg="white").place(x=200, y=620)
            else:
                # print("file1-", "Line-%d" % line_no, file_1_line)
                Label(root, text=f"File1- {line_no} Line- {file_1_line}", font=("Calibri", 15), bg="white").place(x=200, y=620)

            if file_2_line == '':
                # print("file2-", "Line-%d" % line_no, file_2_line)
                Label(root, text=f"File2- {line_no} Line- {file_2_line}", font=("Calibri", 15), bg="white").place(x=200, y=640)
            else:
                # print("file2-", "Line-%d" % line_no, file_2_line)
                Label(root, text=f"File2- {line_no} Line- {file_2_line}", font=("Calibri", 15), bg="white").place(x=200, y=640)

            print()

        file_1_line = file_1.readline()
        file_2_line = file_2.readline()

        line_no += 1

    file_1.close()
    file_2.close()


def get_data():
    choice = choice_entry.get()
    file1_label = Label(root, text="Enter File 1 Name", font=("Calibri", 15), bg="white").place(x=200, y=330)
    file1_entry = Entry(root, width=20, font=("Calibri", 15))
    file1_entry.place(x=200, y=370)

    file2_label = Label(root, text="Enter File 2 Name", font=("Calibri", 15), bg="white").place(x=200, y=430)
    file2_entry = Entry(root, width=20, font=("Calibri", 15))
    file2_entry.place(x=200, y=470)
    if(choice == "PDF"):
        def pdf_check():
            out = pdf2text(file1_entry.get(), file2_entry.get())
            Label(root, text=f"Number of mistakes are : {out}").pack()

        Button(root, text="Enter", command=pdf_check, bg="white", width=15, border=1, font=("Calibri", 15)).place(x=220, y=530)

    elif(choice == "TEXT"):
        def textcheck():
            comparetext(file1_entry.get(), file2_entry.get())

        Button(root, text="Enter", command=textcheck, bg="white", width=15, border=1, font=("Calibri", 15)).place(x=220, y=530)

    elif(choice == "PPT"):
        def pptcheck():
            compare_ppt(file1_entry.get(), file2_entry.get())

        Button(root, text="Enter", command=pptcheck, bg="white", width=15, border=1, font=("Calibri", 15)).place(x=220, y=530)

    elif(choice == 'IMAGE'):
        def imagecheck():
            compare_images(file1_entry.get(), file2_entry.get())
        Button(root, text="Enter", command=imagecheck, bg="white", width=15, border=1, font=("Calibri", 15)).place(x=220, y=530)


Label(root, text="Enter your choice (PPT, TEXT, PDF, IMAGE)", font=("Calibri", 15), bg="white").place(x=120, y=130)


choice_entry = Entry(root, width=20, font=("Calibri", 15))
choice_entry.place(x=200, y=170)

button = Button(root, text="Enter", command=get_data, bg="white", width=15, border=1, font=("Calibri", 15)).place(x=220, y=220)

root.mainloop()
