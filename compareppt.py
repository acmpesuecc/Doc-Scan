import collections 
import collections.abc
from pptx import Presentation

def compare_ppt(p1,p2):
    
    ppt1 = Presentation(p1)
    ppt2 = Presentation(p2)

    text_runs1 = []
    text_runs2 = []

    for slide in ppt1.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs1.append(run.text)
    for slide in ppt2.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs2.append(run.text)
    if collections.Counter(text_runs1) == collections.Counter(text_runs2):
        print("The presentations are same")
    else:
        print("The presentations are different")
        print("The differences are : ")
        print(collections.Counter(text_runs1) - collections.Counter(text_runs2))

